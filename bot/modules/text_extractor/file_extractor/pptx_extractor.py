from ..base_interface import IDocumentTextExtractor
from pptx import Presentation  # python-pptx for PPTX files


class PptxTextExtractor(IDocumentTextExtractor):
    def extract_text(self, file_path: str) -> str:
        if not file_path.endswith('.pptx'):
            raise ValueError("Invalid file type. Expected a .pptx file.")
        text = ""
        try:
            presentation = Presentation(file_path)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            print(f"Error in PPTX extraction: {e}")
        return text
